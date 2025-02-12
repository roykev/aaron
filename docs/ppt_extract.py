import os
import shutil

from pptx import Presentation
import csv
import os
import platform
import subprocess
import sys

# Prerequisites:

# For Windows: pip install python-pptx pillow pywin32
# For Ubuntu: sudo apt-get install libreoffice unoconv
#
# Usage:
#
# Automatic OS detection:
# python script.py /path/to/presentation.pptx
#
# Force specific OS:
# python script.py /path/to/presentation.pptx windows
# python script.py /path/to/presentation.pptx ubuntu

def extract_pdf_slides(pdf_path,output_folder):
    # Convert PDF to images using pdftoppm
    if not shutil.which('pdftoppm'):
        print("Error: pdftoppm not installed. Install with 'sudo apt-get install poppler-utils'")
        return []

    screenshot_command = [
        'pdftoppm',
        '-png',
        pdf_path,
        os.path.join(output_folder, 'slide')
    ]

    subprocess.run(screenshot_command, check=True)
    # Get list of generated screenshot paths
    screenshot_paths = sorted([
        os.path.join(output_folder, f)
        for f in os.listdir(output_folder)
        if f.endswith('.png')
    ])


    return len(screenshot_paths)


def extract_powerpoint_slides(pptx_path, output_folder=None, os_type=None):
    """
    Extract screenshots from PowerPoint slides across different operating systems

    Args:
        pptx_path (str): Path to the PowerPoint presentation
        output_folder (str, optional): Folder to save screenshots
        os_type (str, optional): Override OS detection ('windows' or 'ubuntu')

    Returns:
        list: Paths to saved screenshot images
    """
    # Validate input file
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    # Determine OS if not specified
    if os_type is None:
        current_os = platform.system().lower()
        if current_os == 'windows':
            os_type = 'windows'
        elif current_os == 'linux':
            os_type = 'ubuntu'
        else:
            raise OSError(f"Unsupported operating system: {current_os}")

    # Determine output folder
    if output_folder is None:
        output_folder = os.path.join(
            os.path.dirname(pptx_path),
            f"{os.path.splitext(os.path.basename(pptx_path))[0]}_screenshots"
        )

    # Create output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    # Windows method using win32com
    if os_type == 'windows':
        try:
            import win32com.client

            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            try:
                presentation = powerpoint.Presentations.Open(pptx_path)
                screenshot_paths = []

                for i in range(1, presentation.Slides.Count + 1):
                    # Generate output path for each slide screenshot
                    output_path = os.path.join(output_folder, f"slide_{i}.png")

                    # Export slide as image
                    presentation.Slides[i - 1].Export(output_path, "PNG")
                    screenshot_paths.append(output_path)

                # Close the presentation
                presentation.Close()
            finally:
                # Quit PowerPoint application
                powerpoint.Quit()

            return screenshot_paths

        except ImportError:
            print("Error: pywin32 not installed. Install with 'pip install pywin32'")
            return []

    # Ubuntu method using LibreOffice
    elif os_type == 'ubuntu':
        try:
            # Check if LibreOffice is installed
            if not shutil.which('libreoffice'):
                print("Error: LibreOffice is not installed. Install with 'sudo apt-get install libreoffice'")
                return []

            # Create a temporary directory for conversion
            temp_dir = os.path.join(output_folder, 'temp_converted')
            os.makedirs(temp_dir, exist_ok=True)

            # Get base filename without extension
            base_filename = os.path.splitext(os.path.basename(pptx_path))[0]

            # Convert presentation to PDF first
            conversion_command = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ]

            subprocess.run(conversion_command, check=True)

            # Find the generated PDF file
            pdf_files = [f for f in os.listdir(temp_dir) if f.endswith('.pdf')]
            if not pdf_files:
                print("Error: PDF conversion failed")
                return []

            pdf_path = os.path.join(temp_dir, pdf_files[0])

            # Convert PDF to images using pdftoppm
            if not shutil.which('pdftoppm'):
                print("Error: pdftoppm not installed. Install with 'sudo apt-get install poppler-utils'")
                return []

            screenshot_command = [
                'pdftoppm',
                '-png',
                pdf_path,
                os.path.join(output_folder, 'slide')
            ]

            subprocess.run(screenshot_command, check=True)

            # Get list of generated screenshot paths
            screenshot_paths = sorted([
                os.path.join(output_folder, f)
                for f in os.listdir(output_folder)
                if f.endswith('.png')
            ])

            # Clean up temporary directory
            shutil.rmtree(temp_dir)

            return len(screenshot_paths)

        except subprocess.CalledProcessError as e:
            print(f"Error converting slides: {e}")
            return []

    else:
        raise ValueError(f"Unsupported OS type: {os_type}")


def extract_pptx_content(pptx_path):
    prs = Presentation(pptx_path)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_folder = os.path.join(os.path.dirname(pptx_path),f"{base_name}_extracted")
    os.makedirs(output_folder, exist_ok=True)

    for slide_number, slide in enumerate(prs.slides, start=1):
        # Save text content
        text_file_path = os.path.join(output_folder, f"{slide_number}.txt")
        with open(text_file_path, "w", encoding="utf-8") as text_file:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_file.write(paragraph.text + "\n")

        # Save images
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Placeholder for pictures
                image = shape.image
                image_count += 1
                image_file_path = os.path.join(output_folder, f"{slide_number}_{image_count}.jpg")
                with open(image_file_path, "wb") as img_file:
                    img_file.write(image.blob)

        # Save tables
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_count += 1
                table_file_path = os.path.join(output_folder, f"{slide_number}_{table_count}.csv")
                with open(table_file_path, "w", newline='', encoding="utf-8") as csv_file:
                    csv_writer = csv.writer(csv_file)
                    for row in table.rows:
                        csv_writer.writerow([cell.text for cell in row.cells])

    print(f"Extraction complete. Files saved to '{output_folder}'.")


# Example usage
if __name__ == "__main__":
    # snapshot images
    # Check if a path is provided as an argument
    dirpath = "/home/roy/OneDriver/WORK/ideas/aaron/Miller/AI for business/2024/2/1"
    ppt_path = os.path.join(dirpath,"AI.pptx")
    # if len(sys.argv) < 2:
    #     print("Usage: python script.py /path/to/presentation.pptx [os_type]")
    #     sys.exit(1)
    #
    # ppt_path = sys.argv[1]
    # os_type = sys.argv[2] if len(sys.argv) > 2 else None
    os_type = 'ubuntu'

    try:
        screenshots = extract_powerpoint_slides(ppt_path, os_type=os_type)
        print(f"Screenshots saved. Total slides: {len(screenshots)}")
        for screenshot in screenshots:
            print(f"- {screenshot}")
    except Exception as e:
        print(f"An error occurred: {e}")
# Example usage:
    extract_pptx_content(ppt_path)
